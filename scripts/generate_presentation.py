#!/usr/bin/env python3
"""
Generate PowerPoint presentation from repository analysis.
Uses python-pptx library.
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

def add_title_slide(prs, repo_info):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = repo_info['name']
    subtitle.text = "Architecture & Implementation Guide"

def add_overview_slide(prs, repo_info):
    """Add project overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    tf = content.text_frame
    tf.text = f"Purpose: {repo_info.get('purpose', 'N/A')}"
    
    p = tf.add_paragraph()
    p.text = f"Tech Stack: {repo_info.get('tech_stack', 'N/A')}"

def add_architecture_slide(prs, repo_info):
    """Add architecture slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Architecture"
    tf = content.text_frame
    
    arch_lines = repo_info.get('architecture', 'Architecture details...').split('\n')
    tf.text = arch_lines[0] if arch_lines else ''
    
    for line in arch_lines[1:]:
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 1 if line.startswith('###') else 0

def add_components_slide(prs, repo_info):
    """Add core components slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Components"
    tf = content.text_frame
    
    comp_lines = repo_info.get('components', 'Component details...').split('\n')
    tf.text = comp_lines[0] if comp_lines else ''
    
    for line in comp_lines[1:]:
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()

def add_implementation_slide(prs, repo_info):
    """Add implementation steps slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation Steps"
    tf = content.text_frame
    
    steps_lines = repo_info.get('steps', 'Step-by-step guide...').split('\n')
    tf.text = steps_lines[0] if steps_lines else ''
    
    for line in steps_lines[1:]:
        if line.strip():
            p = tf.add_paragraph()
            p.text = line.strip()

def generate_presentation(repo_info, output_path):
    """Generate PowerPoint presentation."""
    prs = Presentation()
    
    add_title_slide(prs, repo_info)
    add_overview_slide(prs, repo_info)
    add_architecture_slide(prs, repo_info)
    add_components_slide(prs, repo_info)
    add_implementation_slide(prs, repo_info)
    
    prs.save(output_path)
    return output_path

def main():
    if len(sys.argv) < 2:
        print("Usage: generate_presentation.py <repo_info.json> [output.pptx]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else "presentation.pptx"
    
    with open(input_file, 'r', encoding='utf-8') as f:
        repo_info = json.load(f)
    
    pptx_path = generate_presentation(repo_info, output_pptx)
    print(f"✓ Generated PowerPoint: {pptx_path}")

if __name__ == '__main__':
    main()
